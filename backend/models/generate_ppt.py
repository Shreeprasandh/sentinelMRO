import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def apply_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def set_slide_transition_fade(slide):
    """
    Appends a transition XML node directly to the slide element
    to enable an elegant, professional fade animation between slides.
    """
    try:
        from pptx.oxml.xmlchemy import OxmlElement
        # Create <p:transition> element
        transition = OxmlElement('p:transition')
        transition.set('spd', 'med') # Medium speed
        
        # Create <p:fade> element
        fade = OxmlElement('p:fade')
        transition.append(fade)
        
        # Append transition element to the slide root element
        slide._element.append(transition)
    except Exception as e:
        print(f"Could not set slide transition: {e}")

def create_textbox(slide, left, top, width, height, margin=0):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    return tf

def add_slide_header(slide, title_text, category_text="SentinelMRO Core Subsystem"):
    # Header bar background shape
    header_box = slide.shapes.add_shape(
        1, # Rectangle
        Inches(0), Inches(0), Inches(10), Inches(0.9)
    )
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = RGBColor(243, 244, 246) # Light grey (#f3f4f6)
    header_box.line.color.rgb = RGBColor(229, 231, 235) # Very light grey border
    header_box.line.width = Pt(1.0)

    # Title text
    tf = create_textbox(slide, Inches(0.5), Inches(0.12), Inches(6.5), Inches(0.6))
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Arial'
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 58, 138) # Deep Blue (#1e3a8a)

    # Category subtitle
    tf_c = create_textbox(slide, Inches(7.0), Inches(0.25), Inches(2.5), Inches(0.4))
    p_c = tf_c.paragraphs[0]
    p_c.alignment = PP_ALIGN.RIGHT
    p_c.text = category_text
    p_c.font.name = 'Arial'
    p_c.font.size = Pt(9.5)
    p_c.font.color.rgb = RGBColor(13, 148, 136) # Teal (#0d9488)
    p_c.font.bold = True

def add_bullet_point(tf, title, body, bold_title=True):
    p = tf.add_paragraph()
    p.space_before = Pt(4)
    p.space_after = Pt(8)
    
    if bold_title:
        run_title = p.add_run()
        run_title.text = "• " + title + ": "
        run_title.font.bold = True
        run_title.font.size = Pt(12)
        run_title.font.name = 'Arial'
        run_title.font.color.rgb = RGBColor(31, 41, 55) # Dark charcoal (#1f2937)
    else:
        run_title = p.add_run()
        run_title.text = "• " + title
        run_title.font.size = Pt(12)
        run_title.font.name = 'Arial'
        run_title.font.color.rgb = RGBColor(31, 41, 55)

    run_body = p.add_run()
    run_body.text = body
    run_body.font.size = Pt(11)
    run_body.font.name = 'Arial'
    run_body.font.color.rgb = RGBColor(75, 85, 99) # Slate grey (#4b5563)

def build_presentation():
    prs = Presentation()
    
    # Set to widescreen 16:9
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    
    blank_layout = prs.slide_layouts[6]
    
    # White Theme colors
    c_bg = RGBColor(250, 250, 250)    # Soft white background (#fafafa)
    c_blue = RGBColor(30, 58, 138)     # Deep blue primary (#1e3a8a)
    c_teal = RGBColor(13, 148, 136)    # Teal secondary (#0d9488)
    c_card_bg = RGBColor(243, 244, 246) # Light grey card background (#f3f4f6)
    c_border = RGBColor(229, 231, 235)  # Border gray (#e5e7eb)
    
    # Paths for generated images
    current_dir = os.path.dirname(os.path.abspath(__file__))
    img_dir = os.path.abspath(os.path.join(current_dir, "..", "..", "documentation", "images"))
    img_engine = os.path.join(img_dir, "engine_telemetry.png")
    img_fed = os.path.join(img_dir, "federated_learning.png")
    
    # ==================== SLIDE 1: TITLE SLIDE (ELEGANT WHITE) ====================
    slide1 = prs.slides.add_slide(blank_layout)
    apply_background(slide1, c_bg)
    set_slide_transition_fade(slide1)
    
    # Left decorative card line
    accent_box = slide1.shapes.add_shape(
        1, Inches(0.6), Inches(1.4), Inches(0.08), Inches(2.2)
    )
    accent_box.fill.solid()
    accent_box.fill.fore_color.rgb = c_blue
    accent_box.line.fill.background()

    tf1 = create_textbox(slide1, Inches(0.9), Inches(1.3), Inches(8.5), Inches(1.4))
    p_title = tf1.paragraphs[0]
    p_title.text = "SENTINELMRO"
    p_title.font.size = Pt(46)
    p_title.font.bold = True
    p_title.font.name = 'Arial'
    p_title.font.color.rgb = c_blue
    
    p_sub = tf1.add_paragraph()
    p_sub.text = "Decentralized Edge-AI Diagnostics & Immutable Cryptographic Audit Ledger"
    p_sub.font.size = Pt(15.5)
    p_sub.font.name = 'Arial'
    p_sub.font.color.rgb = c_teal
    p_sub.space_before = Pt(6)
    
    # Bottom metadata
    tf_meta = create_textbox(slide1, Inches(0.9), Inches(3.8), Inches(8.0), Inches(1.2))
    p_m1 = tf_meta.paragraphs[0]
    p_m1.text = "Prepared for the Tata Technologies InnoVent-27 Competition (Category 3.2.3.3)"
    p_m1.font.size = Pt(11)
    p_m1.font.bold = True
    p_m1.font.color.rgb = RGBColor(107, 114, 128)
    
    p_m2 = tf_meta.add_paragraph()
    p_m2.text = "Decentralized Prognostics | SQLite Merkle Mountain Range Ledger | Next.js 15 Command Dashboard"
    p_m2.font.size = Pt(9.5)
    p_m2.font.color.rgb = RGBColor(156, 163, 175)
    p_m2.space_before = Pt(4)

    # ==================== SLIDE 2: THE PROBLEM STATEMENT ====================
    slide2 = prs.slides.add_slide(blank_layout)
    apply_background(slide2, c_bg)
    set_slide_transition_fade(slide2)
    add_slide_header(slide2, "The Problem: Centralized Cloud PHM Deficiencies", "Context & Industry Pain Points")
    
    # Left Card
    card_l = slide2.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(4.3), Inches(3.7))
    card_l.fill.solid()
    card_l.fill.fore_color.rgb = c_card_bg
    card_l.line.color.rgb = c_border
    tf_col1 = card_l.text_frame
    tf_col1.word_wrap = True
    tf_col1.margin_left = Inches(0.15)
    tf_col1.margin_right = Inches(0.15)
    tf_col1.margin_top = Inches(0.15)
    
    p_h1 = tf_col1.paragraphs[0]
    p_h1.text = "Legacy Architecture Bottlenecks"
    p_h1.font.size = Pt(15)
    p_h1.font.bold = True
    p_h1.font.color.rgb = c_blue
    p_h1.space_after = Pt(8)
    add_bullet_point(tf_col1, "Bandwidth & Latency", "Uploading gigabytes of high-frequency telemetry streams to clouds causes latency and increases network overhead.")
    add_bullet_point(tf_col1, "Airline Data Privacy", "Airlines are highly protective of operational metrics; central uploading raises severe data security concerns.")

    # Right Card
    card_r = slide2.shapes.add_shape(1, Inches(5.2), Inches(1.3), Inches(4.3), Inches(3.7))
    card_r.fill.solid()
    card_r.fill.fore_color.rgb = c_card_bg
    card_r.line.color.rgb = c_border
    tf_col2 = card_r.text_frame
    tf_col2.word_wrap = True
    tf_col2.margin_left = Inches(0.15)
    tf_col2.margin_right = Inches(0.15)
    tf_col2.margin_top = Inches(0.15)
    
    p_h2 = tf_col2.paragraphs[0]
    p_h2.text = "Audit & Database Vulnerabilities"
    p_h2.font.size = Pt(15)
    p_h2.font.bold = True
    p_h2.font.color.rgb = c_blue
    p_h2.space_after = Pt(8)
    add_bullet_point(tf_col2, "Retroactive Tampering", "Standard maintenance logs stored in centralized databases lack cryptographic seals, risking unauthorized tampering.")
    add_bullet_point(tf_col2, "Non-IID Operational Bias", "Operating conditions vary across MRO hubs. Centralized models become biased toward nominal or high-use fleets.")

    # ==================== SLIDE 3: THE DECENTRALIZED SOLUTION ====================
    slide3 = prs.slides.add_slide(blank_layout)
    apply_background(slide3, c_bg)
    set_slide_transition_fade(slide3)
    add_slide_header(slide3, "The Solution: SentinelMRO Architecture", "Subsystem Overview")
    
    # 3 Column cards layout
    widths = [Inches(2.8), Inches(2.8), Inches(2.8)]
    lefts = [Inches(0.5), Inches(3.6), Inches(6.7)]
    
    titles = [
        "1. Edge-AI Diagnostics",
        "2. Privacy Federated Learning",
        "3. Cryptographic MMR Ledger"
    ]
    
    contents = [
        "Preprocesses 14 sensor channels and regresses RUL on an INT8-quantized TCN model in under 3 milliseconds.",
        "Executes cross-station FedProx local updates and aggregates gradients with Local Differential Privacy (LDP) perturbations.",
        "Validates Ed25519 edge signatures and logs events to an append-only Merkle Mountain Range database in SQLite."
    ]
    
    for i in range(3):
        # Card background
        card = slide3.shapes.add_shape(
            1, lefts[i], Inches(1.3), widths[i], Inches(3.7)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = c_card_bg
        card.line.color.rgb = c_border
        card.line.width = Pt(1.0)
        
        tf_card = card.text_frame
        tf_card.word_wrap = True
        tf_card.margin_left = Inches(0.15)
        tf_card.margin_right = Inches(0.15)
        tf_card.margin_top = Inches(0.15)
        
        p = tf_card.paragraphs[0]
        p.text = titles[i]
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = c_blue
        p.space_after = Pt(8)
        
        p_c = tf_card.add_paragraph()
        p_c.text = contents[i]
        p_c.font.size = Pt(10.5)
        p_c.font.color.rgb = RGBColor(75, 85, 99)

    # ==================== SLIDE 4: PHASE 1 EDGE INFERENCE (WITH IMAGE) ====================
    slide4 = prs.slides.add_slide(blank_layout)
    apply_background(slide4, c_bg)
    set_slide_transition_fade(slide4)
    add_slide_header(slide4, "Phase 1: Quantized Causal TCN Engine", "Edge Inference")
    
    # Left Content Box
    tf_col4 = create_textbox(slide4, Inches(0.5), Inches(1.3), Inches(5.0), Inches(3.7))
    add_bullet_point(tf_col4, "Data Preprocessing", "Strips 7 low-variance constants, keeping 14 features scaled locally via MinMaxScaler [0, 1] over 30 cycles.")
    add_bullet_point(tf_col4, "Dilated Conv1D Network", "Causal padding ensures outputs at cycle t depend only on past cycles. Exponential dilations (1, 2, 4, 8) enable O(1) latency.")
    add_bullet_point(tf_col4, "Quantization & Hot Swap", "Strips PyTorch weight norm hooks, exporting a clean graph. Dynamic INT8 quantization compresses size 4x, achieving < 3ms CPU inference.")
    
    # Right Image
    if os.path.exists(img_engine):
        slide4.shapes.add_picture(img_engine, Inches(5.7), Inches(1.3), width=Inches(3.8), height=Inches(3.7))

    # ==================== SLIDE 5: PHASE 2 FEDERATED LEARNING (WITH IMAGE) ====================
    slide5 = prs.slides.add_slide(blank_layout)
    apply_background(slide5, c_bg)
    set_slide_transition_fade(slide5)
    add_slide_header(slide5, "Phase 2: FedProx & LDP Collaborative Learning", "Federated Optimization")
    
    # Left Content Box
    tf_col5 = create_textbox(slide5, Inches(0.5), Inches(1.3), Inches(5.0), Inches(3.7))
    add_bullet_point(tf_col5, "Non-IID Operational Clusters", "Splits the C-MAPSS dataset into 3 operational clusters: Nominal (cycles <= 60), Stress-State (cycles > 100), and Mixed.")
    add_bullet_point(tf_col5, "FedProx Regularizer Penalty", "Modifies local loss: (mu/2) * ||w - w_t||^2 where mu=0.01. Stabilizes optimization updates under non-IID conditions.")
    add_bullet_point(tf_col5, "Local Differential Privacy", "Adds calibrated Gaussian noise N(0, 0.001^2) to parameters before aggregation to mathematically block model inversion attacks.")
    
    # Right Image
    if os.path.exists(img_fed):
        slide5.shapes.add_picture(img_fed, Inches(5.7), Inches(1.3), width=Inches(3.8), height=Inches(3.7))

    # ==================== SLIDE 6: PHASE 3 CRYPTOGRAPHIC LEDGER ====================
    slide6 = prs.slides.add_slide(blank_layout)
    apply_background(slide6, c_bg)
    set_slide_transition_fade(slide6)
    add_slide_header(slide6, "Phase 3: Cryptographic SQLite MMR Ledger", "Audit & Security")
    
    # Left Column: Principles
    card_l6 = slide6.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(4.3), Inches(3.7))
    card_l6.fill.solid()
    card_l6.fill.fore_color.rgb = c_card_bg
    card_l6.line.color.rgb = c_border
    tf_col6_l = card_l6.text_frame
    tf_col6_l.word_wrap = True
    tf_col6_l.margin_left = Inches(0.15)
    tf_col6_l.margin_top = Inches(0.15)
    
    p_hl6 = tf_col6_l.paragraphs[0]
    p_hl6.text = "Cryptographic Ledger Principles"
    p_hl6.font.size = Pt(14)
    p_hl6.font.bold = True
    p_hl6.font.color.rgb = c_blue
    p_hl6.space_after = Pt(8)
    
    add_bullet_point(tf_col6_l, "Ed25519 Gateway Handshake", "Every edge submission is signed by the station's private key. Gateway checks signatures in headers using public keys.")
    add_bullet_point(tf_col6_l, "Merkle Mountain Range (MMR)", "Appends nodes to a binary tree in post-order flat indexing. Enables O(log n) inclusion proofs and rapid integrity verification checks.")

    # Right Column: Backdoor Tamper Schema
    card_r6 = slide6.shapes.add_shape(1, Inches(5.2), Inches(1.3), Inches(4.3), Inches(3.7))
    card_r6.fill.solid()
    card_r6.fill.fore_color.rgb = c_card_bg
    card_r6.line.color.rgb = c_border
    tf_col6_r = card_r6.text_frame
    tf_col6_r.word_wrap = True
    tf_col6_r.margin_left = Inches(0.15)
    tf_col6_r.margin_top = Inches(0.15)
    
    p_hr6 = tf_col6_r.paragraphs[0]
    p_hr6.text = "Tamper Backdoor Demonstration"
    p_hr6.font.size = Pt(14)
    p_hr6.font.bold = True
    p_hr6.font.color.rgb = c_blue
    p_hr6.space_after = Pt(8)
    
    add_bullet_point(tf_col6_r, "Injecting Corrupted Payload", "Updates SQLite columns directly (e.g. component_id='TAMPERED-ENG') but leaves the MMR nodes unchanged.")
    add_bullet_point(tf_col6_r, "Verification Alert", "Recomputing leaf hashes during audits reveals a discrepancy with stored node hashes. The UI validation alarm instantly flashes red.")

    # ==================== SLIDE 7: PHASE 4 COMMAND DASHBOARD ====================
    slide7 = prs.slides.add_slide(blank_layout)
    apply_background(slide7, c_bg)
    set_slide_transition_fade(slide7)
    add_slide_header(slide7, "Phase 4: Next.js 15 Command Console", "User Interface & Control")
    
    tf_col7 = create_textbox(slide7, Inches(0.5), Inches(1.3), Inches(9.0), Inches(3.7))
    
    add_bullet_point(tf_col7, "Premium Industrial Theme", "Designed on zinc-950/zinc-50 with responsive sidebar navigation, modern typography ( Outfit/Inter), and micro-animations.")
    add_bullet_point(tf_col7, "Real-time Telemetry curves", "Visualizes TCN RUL outputs alongside core speed temperature and bypass ratio trends using high-frequency Recharts curves.")
    add_bullet_point(tf_col7, "Auditing Control Console", "Provides buttons to Verify Ecosystem Integrity or Simulate Malicious Database Injection. Alerts flash green on success and red on mismatch.")

    # ==================== SLIDE 8: WORKED TRACE EXAMPLE ====================
    slide8 = prs.slides.add_slide(blank_layout)
    apply_background(slide8, c_bg)
    set_slide_transition_fade(slide8)
    add_slide_header(slide8, "Worked Trace Example: Maintenance Logging", "Ecosystem in Action")
    
    # Left Card
    card_l8 = slide8.shapes.add_shape(1, Inches(0.5), Inches(1.3), Inches(4.3), Inches(3.7))
    card_l8.fill.solid()
    card_l8.fill.fore_color.rgb = c_card_bg
    card_l8.line.color.rgb = c_border
    tf_l8 = card_l8.text_frame
    tf_l8.word_wrap = True
    tf_l8.margin_left = Inches(0.15)
    tf_l8.margin_top = Inches(0.15)
    
    p_hl8 = tf_l8.paragraphs[0]
    p_hl8.text = "1. Logging & Signing"
    p_hl8.font.size = Pt(14)
    p_hl8.font.bold = True
    p_hl8.font.color.rgb = c_blue
    p_hl8.space_after = Pt(8)
    
    add_bullet_point(tf_l8, "Generate Payload", "Edge node compiles: Timestamp|Component_ID|Action|Tech_ID|Health.")
    add_bullet_point(tf_l8, "Ed25519 Signature", "Signs payload using private key. Submitted in header to API gateway.")
    add_bullet_point(tf_l8, "Merkle Leaf Append", "Gateway verifies signature, writes row, and appends leaf hash: '5f8a0dc4f67c...' to SQLite.")

    # Right Card
    card_r8 = slide8.shapes.add_shape(1, Inches(5.2), Inches(1.3), Inches(4.3), Inches(3.7))
    card_r8.fill.solid()
    card_r8.fill.fore_color.rgb = c_card_bg
    card_r8.line.color.rgb = c_border
    tf_r8 = card_r8.text_frame
    tf_r8.word_wrap = True
    tf_r8.margin_left = Inches(0.15)
    tf_r8.margin_top = Inches(0.15)
    
    p_hr8 = tf_r8.paragraphs[0]
    p_hr8.text = "2. Tampering & Alerting"
    p_hr8.font.size = Pt(14)
    p_hr8.font.bold = True
    p_hr8.font.color.rgb = c_blue
    p_hr8.space_after = Pt(8)
    
    add_bullet_point(tf_r8, "Malicious Modification", "Intruder updates SQLite directly (e.g. altering the component_id to 'TAMPERED-ENG').")
    add_bullet_point(tf_r8, "Hash Discrepancy", "Auditor re-computes Leaf 1 hash, producing '7a11c8...', which mismatches the stored '5f8a0dc4f67c...' value.")
    add_bullet_point(tf_r8, "Ecosystem Integrity Alarm", "Mismatch invalidates the MMR root. The command console validation indicator flashes red and displays an alert.")

    # ==================== SLIDE 9: STATUS & CHECKLIST ====================
    slide9 = prs.slides.add_slide(blank_layout)
    apply_background(slide9, c_bg)
    set_slide_transition_fade(slide9)
    add_slide_header(slide9, "Conclusion: SentinelMRO InnoVent-27 Ready", "Competitive Summary")
    
    tf_col9 = create_textbox(slide9, Inches(0.5), Inches(1.3), Inches(9.0), Inches(3.7))
    
    add_bullet_point(tf_col9, "Edge AI Latency", "Model compiled and quantized dynamically, performing edge prediction on CPU in under 3ms.")
    add_bullet_point(tf_col9, "Collaborative FedProx", "Federated loop stabilizes updates under non-IID conditions while LDP noise aggregation secures gradients.")
    add_bullet_point(tf_col9, "Cryptographic Trust", "SQLite-backed MMR append, verify, and backdoor tampering functions fully complete.")
    add_bullet_point(tf_col9, "Next.js 15 command console", "Command console provides high-fidelity, real-time visual control room ready for live demonstration.")

    # Save
    ppt_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "documentation"))
    os.makedirs(ppt_dir, exist_ok=True)
    ppt_path = os.path.join(ppt_dir, "SentinelMRO_Project_Presentation.pptx")
    prs.save(ppt_path)
    print(f"Presentation saved successfully at: {ppt_path}")

if __name__ == "__main__":
    build_presentation()
